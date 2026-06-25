from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── PALETA TRADICIONALISMO GAÚCHO ───────────────────────────────────────────
VERMELHO    = RGBColor(0x8B, 0x1A, 0x1A)   # vermelho escuro (bombachas)
VERDE       = RGBColor(0x1A, 0x4A, 0x2A)   # verde RS
DOURADO     = RGBColor(0xC8, 0x96, 0x0C)   # dourado / amarelo RS
CREME       = RGBColor(0xF5, 0xF0, 0xE8)   # pergaminho
MARROM      = RGBColor(0x5C, 0x3A, 0x1E)   # couro
BRANCO      = RGBColor(0xFF, 0xFF, 0xFF)
CINZA       = RGBColor(0x44, 0x44, 0x44)
CINZA_CLARO = RGBColor(0xF0, 0xED, 0xE8)
VERDE_CLARO = RGBColor(0xD6, 0xE8, 0xD8)
VERM_CLARO  = RGBColor(0xF2, 0xD6, 0xD6)
DOUR_CLARO  = RGBColor(0xFD, 0xF0, 0xC8)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# ── HELPERS ──────────────────────────────────────────────────────────────────
def rect(sl, l, t, w, h, fill, line=None):
    s = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line
    else: s.line.fill.background()
    return s

def txt(sl, text, l, t, w, h, size, bold=False, color=CINZA,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tx = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tx.word_wrap = wrap
    tf = tx.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tx

def header(sl, title, subtitle=None, bg=VERMELHO, accent=DOURADO):
    rect(sl, 0, 0, 13.33, 1.35, bg)
    rect(sl, 0, 1.35, 13.33, 0.07, accent)   # linha dourada
    txt(sl, title, 0.45, 0.18, 12.0, 0.85, 26, bold=True, color=BRANCO)
    if subtitle:
        txt(sl, subtitle, 0.45, 0.92, 12.0, 0.38, 12, color=DOURADO, italic=True)
    rect(sl, 0, 1.42, 13.33, 6.08, CREME)

def stat_block(sl, l, t, w, h, value, label, bg=VERDE, val_color=BRANCO, lbl_color=BRANCO):
    rect(sl, l, t, w, h, bg)
    txt(sl, value, l+0.1, t+0.18, w-0.2, h*0.55, 28, bold=True,
        color=val_color, align=PP_ALIGN.CENTER)
    txt(sl, label, l+0.1, t+h*0.58, w-0.2, h*0.42, 11,
        color=lbl_color, align=PP_ALIGN.CENTER, wrap=True)

def card(sl, l, t, w, h, title, body, bg=VERDE_CLARO, title_bg=VERDE,
         title_color=BRANCO, body_color=CINZA, title_size=13, body_size=11):
    rect(sl, l, t, w, 0.42, title_bg)
    txt(sl, title, l+0.12, t+0.06, w-0.24, 0.34, title_size,
        bold=True, color=title_color, align=PP_ALIGN.LEFT)
    rect(sl, l, t+0.42, w, h-0.42, bg, line=RGBColor(0xCC,0xCC,0xCC))
    txt(sl, body, l+0.15, t+0.5, w-0.3, h-0.6, body_size,
        color=body_color, wrap=True)

def step_block(sl, l, t, w, h, num, title, body, bg=VERMELHO):
    rect(sl, l, t, 0.5, h, bg)
    txt(sl, num, l+0.05, t+h/2-0.28, 0.4, 0.55, 22,
        bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    rect(sl, l+0.5, t, w-0.5, h, CREME, line=RGBColor(0xCC,0xBB,0xAA))
    txt(sl, title, l+0.65, t+0.1, w-0.8, 0.35, 12, bold=True, color=VERMELHO)
    txt(sl, body, l+0.65, t+0.42, w-0.8, h-0.52, 10.5, color=CINZA, wrap=True)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — CAPA
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
# Fundo verde escuro
rect(s, 0, 0, 13.33, 7.5, VERDE)
# Faixa vermelha central
rect(s, 0, 2.55, 13.33, 2.55, VERMELHO)
# Linhas douradas
rect(s, 0, 2.45, 13.33, 0.12, DOURADO)
rect(s, 0, 5.08, 13.33, 0.12, DOURADO)
# Instituição
txt(s, "ANTONIO MENEGHETTI FACULDADE  |  Sistemas de Informação",
    0.5, 0.3, 12.3, 0.55, 13, color=DOURADO, align=PP_ALIGN.CENTER)
txt(s, "Projeto de Trabalho de Conclusão de Curso",
    0.5, 0.78, 12.3, 0.45, 11, italic=True,
    color=RGBColor(0xCC,0xDD,0xBB), align=PP_ALIGN.CENTER)
# Título principal
txt(s, "Desenvolvimento e Avaliação de um\nEcossistema Digital para Gestão de\nAssociados e Eventos em Entidades\nTradicionalistasGaúchas",
    0.8, 2.65, 11.7, 2.3, 22, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
# Rodapé
txt(s, "Yuri Garcia Baptista",
    0.5, 5.45, 12.3, 0.5, 15, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
txt(s, "Orientador: Prof. Dr. Felipe Becker Nunes",
    0.5, 5.92, 12.3, 0.4, 12, color=DOURADO, align=PP_ALIGN.CENTER)
txt(s, "Restinga Seca, RS  ·  Junho 2026",
    0.5, 6.75, 12.3, 0.4, 11, italic=True,
    color=RGBColor(0xAA,0xCC,0xAA), align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — CONTEXTUALIZAÇÃO
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header(s, "Contextualização", "O Movimento Tradicionalista Gaúcho e o desafio da gestão")

# 3 grandes stats
stat_block(s, 0.4, 1.7, 3.8, 1.8, "+1.700", "entidades\nfiliadas no RS", bg=VERDE)
stat_block(s, 4.75, 1.7, 3.8, 1.8, "100 mil", "eventos\npor ano", bg=VERMELHO)
stat_block(s, 9.1, 1.7, 3.8, 1.8, "R$ 1,5 bi", "movimentados\npor ano", bg=MARROM)

# Linha separadora
rect(s, 0.4, 3.7, 12.53, 0.04, DOURADO)

txt(s, "Associados, invernadas artísticas, bailes e fandangos sustentam uma identidade cultural única no Brasil.",
    0.5, 3.85, 12.3, 0.5, 13, color=CINZA, align=PP_ALIGN.CENTER, italic=True)

# 2 cards problema
card(s, 0.4, 4.55, 5.9, 2.55,
     "Gestão ainda manual",
     "Planilhas, fichas físicas e comunicação informal.\nRetrabalho, inconsistências e baixa rastreabilidade.\nDificuldade na continuidade com troca de diretoria.",
     bg=VERM_CLARO, title_bg=VERMELHO)

card(s, 6.8, 4.55, 5.93, 2.55,
     "Lacuna tecnológica",
     "Não existe sistema integrado específico para\nentidades tradicionalistas gaúchas.\nGestores voluntários sobrecarregados.",
     bg=DOUR_CLARO, title_bg=MARROM)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — QUESTÃO E OBJETIVOS
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header(s, "Objetivos da Pesquisa")

# Questão central
rect(s, 0.4, 1.6, 12.53, 1.25, VERDE)
rect(s, 0.4, 1.6, 0.12, 1.25, DOURADO)
txt(s, "Como um ecossistema digital integrado pode qualificar a gestão de associados\ne eventos sociais em entidades tradicionalistas gaúchas?",
    0.65, 1.72, 12.1, 1.0, 14, color=BRANCO, italic=True, align=PP_ALIGN.LEFT)

# Objetivo geral
rect(s, 0.4, 3.05, 12.53, 0.65, CREME, line=DOURADO)
txt(s, "Objetivo Geral:", 0.6, 3.12, 3.0, 0.35, 12, bold=True, color=VERMELHO)
txt(s, "Desenvolver e avaliar um ecossistema digital entidade-associado para o contexto das entidades tradicionalistas gaúchas.",
    3.4, 3.12, 9.3, 0.5, 12, color=CINZA)

# 4 objetivos específicos em blocos
for i, (letra, texto) in enumerate([
    ("a", "Identificar requisitos de gestão de associados, mensalidades e eventos"),
    ("b", "Descrever a arquitetura e componentes do ecossistema proposto"),
    ("c", "Implementar as funcionalidades básicas do MVP"),
    ("d", "Avaliar o sistema com usuários reais do CPF Pia do Sul"),
]):
    col = i % 2; row = i // 2
    lx = 0.4 + col * 6.5; ty = 3.9 + row * 1.52
    rect(s, lx, ty, 6.1, 1.35, CREME, line=RGBColor(0xBB,0xAA,0x99))
    rect(s, lx, ty, 0.55, 1.35, VERDE)
    txt(s, letra+")", lx+0.05, ty+0.38, 0.45, 0.55, 18, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    txt(s, texto, lx+0.7, ty+0.25, 5.2, 0.85, 11.5, color=CINZA, wrap=True)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — METODOLOGIA
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header(s, "Metodologia", "Pesquisa aplicada e qualitativa  ·  Estudo de caso: CPF Pia do Sul, Santa Maria/RS")

# 3 etapas em blocos horizontais
step_block(s, 0.4, 1.6, 4.0, 1.75, "1",
           "Levantamento de Requisitos",
           "Entrevistas semiestruturadas\nAnálise documental\nObservação dos processos atuais",
           bg=VERDE)
step_block(s, 4.7, 1.6, 4.0, 1.75, "2",
           "Desenvolvimento do MVP",
           "Abordagem DSDM\nTimeboxing + priorização MoSCoW\nCiclos curtos com validação",
           bg=VERMELHO)
step_block(s, 9.0, 1.6, 4.0, 1.75, "3",
           "Avaliação com Usuários",
           "System Usability Scale (SUS)\nEntrevistas semiestruturadas\nObservação de cenários de uso",
           bg=MARROM)

# Setas entre etapas
txt(s, "▶", 4.3, 2.3, 0.45, 0.55, 18, color=DOURADO, align=PP_ALIGN.CENTER)
txt(s, "▶", 8.6, 2.3, 0.45, 0.55, 18, color=DOURADO, align=PP_ALIGN.CENTER)

# Por que DSDM?
rect(s, 0.4, 3.55, 12.53, 0.06, DOURADO)
txt(s, "Por que DSDM?", 0.5, 3.75, 3.5, 0.4, 13, bold=True, color=VERMELHO)
txt(s, "O TCC tem prazo de entrega fixo e escopo variável — exatamente o cenário para o qual o DSDM foi projetado.",
    0.5, 4.15, 6.2, 0.6, 11.5, color=CINZA, wrap=True)

# MoSCoW visual
mx = 7.0
for i, (label, desc, bg) in enumerate([
    ("Must",   "Essencial para o MVP", VERDE),
    ("Should", "Importante, negociável", MARROM),
    ("Could",  "Desejável se possível", RGBColor(0x9B,0x59,0x10)),
    ("Won't",  "Fora do escopo", RGBColor(0x88,0x88,0x88)),
]):
    rect(s, mx + i*1.52, 3.75, 1.45, 2.85, bg)
    txt(s, label, mx + i*1.52 + 0.08, 3.82, 1.3, 0.45, 13, bold=True,
        color=BRANCO, align=PP_ALIGN.CENTER)
    txt(s, desc, mx + i*1.52 + 0.08, 4.32, 1.3, 2.1, 10,
        color=BRANCO, align=PP_ALIGN.CENTER, wrap=True)

# Participantes
rect(s, 0.4, 6.55, 12.53, 0.7, VERDE)
txt(s, "Participantes da avaliação:  6 a 10 pessoas  —  membros da diretoria, responsáveis por eventos e associados",
    0.6, 6.66, 12.1, 0.45, 12, color=BRANCO, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — ARQUITETURA
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header(s, "Arquitetura do Ecossistema Digital",
       "Hexagonal Architecture + Clean Architecture  ·  API RESTful  ·  TypeScript end-to-end")

def abox(sl, l, t, w, h, bg, title, sub="", border=None):
    rect(sl, l, t, w, h, bg, line=border)
    txt(sl, title, l+0.1, t+0.1, w-0.2, 0.4, 12, bold=True,
        color=BRANCO, align=PP_ALIGN.CENTER)
    if sub:
        txt(sl, sub, l+0.1, t+0.48, w-0.2, h-0.58, 9.5,
            color=RGBColor(0xEE,0xEE,0xEE), align=PP_ALIGN.CENTER, wrap=True)

# CLIENTES
abox(s, 0.5, 1.6, 3.5, 1.1, VERDE,
     "Next.js", "Painel Web\n(Administração)")
abox(s, 9.3, 1.6, 3.5, 1.1, VERDE,
     "React Native + Expo", "App Mobile\n(Associado)")

# Setas HTTP
txt(s, "HTTP / REST ▼", 3.9, 1.85, 2.0, 0.4, 10, color=CINZA, align=PP_ALIGN.CENTER, italic=True)
txt(s, "▼ HTTP / REST", 7.4, 1.85, 2.0, 0.4, 10, color=CINZA, align=PP_ALIGN.CENTER, italic=True)

# API NESTJS
abox(s, 3.4, 2.85, 6.5, 1.0, VERMELHO,
     "NestJS — API REST",
     "Controllers  ·  Guards  ·  DTOs  ·  Swagger")

txt(s, "▼", 6.4, 3.95, 0.55, 0.4, 16, color=DOURADO, align=PP_ALIGN.CENTER, bold=True)

# CAMADA APLICAÇÃO
abox(s, 3.4, 4.4, 6.5, 0.85, MARROM,
     "Camada de Aplicação",
     "Use Cases  ·  Ports  ·  Orquestração de fluxos")

txt(s, "▼", 6.4, 5.3, 0.55, 0.4, 16, color=DOURADO, align=PP_ALIGN.CENTER, bold=True)

# DOMÍNIO
rect(s, 3.4, 5.72, 6.5, 0.85, RGBColor(0x3A, 0x22, 0x0E))
rect(s, 3.4, 5.72, 0.12, 0.85, DOURADO)
rect(s, 9.78, 5.72, 0.12, 0.85, DOURADO)
txt(s, "Camada de Domínio", 3.55, 5.76, 6.2, 0.38, 12, bold=True, color=DOURADO, align=PP_ALIGN.CENTER)
txt(s, "Entidades  ·  Regras de Negócio  ·  Invariantes",
    3.55, 6.1, 6.2, 0.35, 9.5, color=CREME, align=PP_ALIGN.CENTER)

# INFRA
abox(s, 0.5, 5.72, 2.6, 1.45, RGBColor(0x44,0x44,0x44),
     "PostgreSQL", "via TypeORM\nAdapter")
abox(s, 10.2, 5.72, 2.6, 1.45, RGBColor(0x44,0x44,0x44),
     "Serviços Externos", "Gateways de Pagamento\nAdapter")

# Setas domínio → infra
txt(s, "▼", 1.75, 5.2, 0.55, 0.5, 14, color=CINZA, align=PP_ALIGN.CENTER)
txt(s, "▼", 11.3, 5.2, 0.55, 0.5, 14, color=CINZA, align=PP_ALIGN.CENTER)

# Nota rodapé
rect(s, 0, 7.15, 13.33, 0.35, VERDE)
txt(s, "TypeScript unificado em todas as camadas  ·  JWT + RBAC  ·  API documentada via Swagger",
    0.5, 7.19, 12.3, 0.28, 10.5, color=BRANCO, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — STACK TECNOLÓGICA
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header(s, "Stack Tecnológica")

camadas = [
    ("Backend", [("NestJS", "Framework Node.js, arquitetura modular"), ("TypeORM + PostgreSQL", "ORM + BD relacional, migrations"), ("JWT + RBAC", "Autenticação e controle de acesso")], VERMELHO),
    ("Frontend Web", [("Next.js", "SSR/SSG, roteamento por arquivo"), ("React Query", "Cache e estado de dados remotos"), ("Tailwind + shadcn/ui", "Estilização e componentes acessíveis")], VERDE),
    ("Mobile", [("React Native", "Código nativo multiplataforma"), ("Expo", "Android e iOS de uma só base"), ("TypeScript", "Compartilhamento de tipos com backend")], MARROM),
]

for ci, (titulo, items, cor) in enumerate(camadas):
    lx = 0.4 + ci * 4.32
    rect(s, lx, 1.6, 4.1, 0.55, cor)
    txt(s, titulo, lx+0.1, 1.65, 3.9, 0.45, 15, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    for ri, (nome, desc) in enumerate(items):
        ty = 2.3 + ri * 1.55
        bg = CREME if ri % 2 == 0 else CINZA_CLARO
        rect(s, lx, ty, 4.1, 1.45, bg, line=cor)
        rect(s, lx, ty, 0.1, 1.45, cor)
        txt(s, nome, lx+0.22, ty+0.12, 3.75, 0.42, 13, bold=True, color=cor)
        txt(s, desc, lx+0.22, ty+0.55, 3.75, 0.8, 10.5, color=CINZA, wrap=True)

# Linha TypeScript unificado
rect(s, 0.4, 7.0, 12.53, 0.28, VERDE)
txt(s, "TypeScript end-to-end — compartilhamento de tipos, contratos e consistência entre todas as camadas",
    0.5, 7.03, 12.2, 0.24, 10.5, color=BRANCO, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — REQUISITOS
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header(s, "Requisitos Funcionais — Priorização MoSCoW")

modulos = [
    ("Gestão de\nAssociados", [
        ("MUST", "Cadastro, edição e status do associado"),
        ("MUST", "Consulta e filtros de associados"),
        ("SHOULD", "Registro de dependentes"),
    ], VERDE, VERDE_CLARO),
    ("Controle de\nMensalidades", [
        ("MUST", "Registro de pagamentos"),
        ("MUST", "Histórico por associado"),
        ("SHOULD", "Relatório de inadimplência"),
    ], VERMELHO, VERM_CLARO),
    ("Eventos e\nReservas", [
        ("MUST", "Cadastro de bailes e fandangos"),
        ("MUST", "Reserva de mesa (fluxo mediado)"),
        ("MUST", "App: consultar reservas e eventos"),
    ], MARROM, DOUR_CLARO),
]

for ci, (mod, reqs, cor, bg) in enumerate(modulos):
    lx = 0.4 + ci * 4.32
    rect(s, lx, 1.6, 4.1, 0.78, cor)
    txt(s, mod, lx+0.1, 1.65, 3.9, 0.7, 14, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    rect(s, lx, 2.38, 4.1, 4.75, bg, line=cor)
    for ri, (prio, desc) in enumerate(reqs):
        ty = 2.55 + ri * 1.45
        prio_cor = VERDE if prio == "MUST" else MARROM
        rect(s, lx+0.18, ty, 0.85, 0.32, prio_cor)
        txt(s, prio, lx+0.18, ty+0.03, 0.85, 0.28, 9, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
        txt(s, desc, lx+0.18, ty+0.38, 3.6, 0.9, 11, color=CINZA, wrap=True)

# Won't have
rect(s, 0.4, 7.05, 12.53, 0.3, RGBColor(0x88,0x88,0x88))
txt(s, "Won't Have:  integração CIT/MTG  ·  pagamento direto no app  ·  rodeios artísticos",
    0.6, 7.09, 12.1, 0.24, 10.5, color=BRANCO, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — RESULTADOS ESPERADOS
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header(s, "Resultados Esperados", "Entregas do MVP + Critérios de sucesso")

# Lado esquerdo: MVP
rect(s, 0.4, 1.6, 0.08, 5.6, DOURADO)
txt(s, "ENTREGAS DO MVP", 0.6, 1.65, 5.8, 0.4, 13, bold=True, color=VERMELHO)

itens_mvp = [
    ("Painel Web (Next.js)", "Associados · Mensalidades · Eventos · Reservas"),
    ("App Mobile (Expo)", "Autenticação · Dados · Reservas · Eventos"),
    ("API REST (NestJS)", "Casos de uso · JWT · Swagger · TypeORM"),
]
for i, (nome, desc) in enumerate(itens_mvp):
    ty = 2.2 + i * 1.6
    rect(s, 0.55, ty, 5.65, 1.45, CREME, line=DOURADO)
    rect(s, 0.55, ty, 0.08, 1.45, VERDE)
    txt(s, nome, 0.75, ty+0.14, 5.25, 0.42, 13, bold=True, color=VERDE)
    txt(s, desc, 0.75, ty+0.6, 5.25, 0.7, 11, color=CINZA)

# Divisória
rect(s, 6.6, 1.6, 0.06, 5.6, DOURADO)

# Lado direito: critérios
txt(s, "CRITÉRIOS DE SUCESSO", 6.85, 1.65, 6.1, 0.4, 13, bold=True, color=VERMELHO)

criterios = [
    "Redução do retrabalho administrativo",
    "Facilidade de uso (SUS + entrevistas)",
    "Clareza das informações nas interfaces",
    "Rapidez para consultar mensalidades e reservas",
    "Aceitação pelos usuários — incluindo\nperfis com menor maturidade digital",
]
for i, c in enumerate(criterios):
    ty = 2.25 + i * 1.02
    rect(s, 6.85, ty, 0.38, 0.38, VERDE)
    txt(s, "✓", 6.85, ty+0.01, 0.38, 0.36, 14, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    txt(s, c, 7.32, ty+0.02, 5.7, 0.85, 11.5, color=CINZA, wrap=True)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — CRONOGRAMA
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header(s, "Cronograma — TCC 2026/2", "Roadmap de entregas incrementais alinhado ao DSDM")

meses = ["Jul", "Ago", "Set", "Out", "Nov"]
ativs = [
    ("Levantamento de requisitos",           [1,0,0,0,0], VERDE),
    ("Modelagem do banco de dados",          [1,0,0,0,0], VERDE),
    ("Dev: Associados + Mensalidades",       [1,1,0,0,0], VERMELHO),
    ("Dev: Eventos + Reservas + Ingressos",  [0,1,1,0,0], VERMELHO),
    ("Dev: App Mobile",                      [0,0,1,0,0], VERMELHO),
    ("Testes funcionais e integração",       [0,0,1,1,0], MARROM),
    ("Avaliação com usuários",               [0,0,0,1,0], MARROM),
    ("Escrita da monografia",                [1,1,1,1,1], RGBColor(0x33,0x66,0x99)),
    ("Entrega da monografia",                [0,0,0,0,1], DOURADO),
]

col_w = 1.48
row_h = 0.52
x0, y0 = 0.35, 1.55
atv_w = 5.0

# Cabeçalho meses
for i, m in enumerate(meses):
    rect(s, x0+atv_w+i*col_w, y0, col_w-0.04, 0.42, VERDE)
    txt(s, m, x0+atv_w+i*col_w, y0+0.02, col_w-0.04, 0.38, 13,
        bold=True, color=BRANCO, align=PP_ALIGN.CENTER)

for ri, (ativ, marks, cor) in enumerate(ativs):
    yy = y0+0.46+ri*row_h
    bg = CREME if ri%2==0 else CINZA_CLARO
    rect(s, x0, yy, atv_w-0.04, row_h-0.04, bg)
    txt(s, ativ, x0+0.12, yy+0.08, atv_w-0.22, row_h-0.12, 11, color=CINZA)
    for ci, m in enumerate(marks):
        rect(s, x0+atv_w+ci*col_w, yy, col_w-0.04, row_h-0.04, bg)
        if m:
            rect(s, x0+atv_w+ci*col_w+0.1, yy+0.08, col_w-0.24, row_h-0.2, cor)

# Milestones
rect(s, 0.35, 6.68, 12.58, 0.65, VERDE)
txt(s, "Milestones:  Requisitos + Modelagem até Jul/26  ·  MVP completo até Set/26  ·  Avaliação até Out/26  ·  Entrega Nov/26",
    0.5, 6.8, 12.2, 0.38, 11, color=BRANCO, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — ENCERRAMENTO
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
rect(s, 0, 0, 13.33, 7.5, VERDE)
rect(s, 0, 0, 13.33, 0.12, DOURADO)
rect(s, 0, 7.38, 13.33, 0.12, DOURADO)
rect(s, 0, 3.0, 13.33, 0.1, DOURADO)

txt(s, "Obrigado!", 0.5, 0.5, 12.3, 1.1, 44, bold=True,
    color=BRANCO, align=PP_ALIGN.CENTER)
txt(s, "Yuri Garcia Baptista  ·  yurigamergarcia@gmail.com",
    0.5, 1.5, 12.3, 0.5, 13, color=DOURADO, align=PP_ALIGN.CENTER)
txt(s, "Orientador: Prof. Dr. Felipe Becker Nunes",
    0.5, 2.1, 12.3, 0.4, 12, italic=True,
    color=RGBColor(0xCC,0xDD,0xBB), align=PP_ALIGN.CENTER)

pontos = [
    ("Proposta aplicada", "Parceria real com o CPF Pia do Sul, Santa Maria/RS"),
    ("Arquitetura sólida", "Hexagonal + Clean Architecture + stack TypeScript moderna"),
    ("Metodologia clara", "3 etapas: requisitos → MVP → avaliação com usuários reais"),
    ("Contribuição", "Modernização administrativa sem perda da identidade cultural"),
]
for i, (titulo, desc) in enumerate(pontos):
    col = i % 2; row = i // 2
    lx = 0.5 + col * 6.5; ty = 3.3 + row * 1.8
    rect(s, lx, ty, 6.1, 1.55, RGBColor(0x0F, 0x2E, 0x18))
    rect(s, lx, ty, 0.12, 1.55, DOURADO)
    txt(s, titulo, lx+0.3, ty+0.2, 5.6, 0.42, 14, bold=True, color=DOURADO)
    txt(s, desc, lx+0.3, ty+0.65, 5.6, 0.75, 11.5, color=CREME, wrap=True)

prs.save(r"e:\TCC\docs\apresentacao_v2.pptx")
print("OK — 10 slides gerados")
